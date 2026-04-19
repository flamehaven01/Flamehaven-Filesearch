"""
Format Backend Plugin system for Flamehaven FileSearch.

Architecture absorbed from Docling (backend/abstract_backend.py):
  AbstractDocumentBackend -> supports() + extract()

Each backend encapsulates one format family's extraction logic.
The BackendRegistry maps file extensions to their backend class.
New formats are registered via BackendRegistry.register() without
touching the dispatcher.

Zero new external dependencies — optional heavy libs (pymupdf, python-docx,
etc.) are still guarded with ImportError inside each backend.
"""

import logging
import subprocess
from abc import ABC, abstractmethod
from typing import Dict, List, Optional, Set, Type

from .format_parsers import (
    extract_csv,
    extract_html,
    extract_image,
    extract_latex,
    extract_vtt,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Abstract base
# ---------------------------------------------------------------------------


class AbstractFormatBackend(ABC):
    """Base class for all format-specific text extractors.

    Subclasses implement extract() and declare their supported extensions
    via the supported_extensions class variable.
    """

    supported_extensions: Set[str] = set()

    @abstractmethod
    def extract(self, file_path: str) -> str:
        """Return extracted plain-text content, or empty string on failure."""

    def _warn_missing(self, lib: str, extra: str) -> None:
        logger.warning(
            "[%s] %s not installed. Run: pip install flamehaven-filesearch[%s]",
            self.__class__.__name__,
            lib,
            extra,
        )

    def _read_plain(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                return fh.read()
        except OSError:
            return ""


# ---------------------------------------------------------------------------
# Backend registry
# ---------------------------------------------------------------------------


class BackendRegistry:
    """Maps file extensions to their AbstractFormatBackend class.

    Usage::

        registry = BackendRegistry.default()
        backend = registry.get(".pdf")
        text = backend.extract(file_path)
    """

    def __init__(self) -> None:
        self._map: Dict[str, Type[AbstractFormatBackend]] = {}

    def register(self, backend_cls: Type[AbstractFormatBackend]) -> None:
        """Register a backend class for all its declared extensions."""
        for ext in backend_cls.supported_extensions:
            self._map[ext.lower()] = backend_cls

    def get(self, ext: str) -> Optional[Type[AbstractFormatBackend]]:
        """Return the backend class for a given extension, or None."""
        return self._map.get(ext.lower())

    def supported_extensions(self) -> Set[str]:
        return set(self._map.keys())

    @classmethod
    def default(cls) -> "BackendRegistry":
        """Build and return the canonical backend registry."""
        registry = cls()
        for backend_cls in _ALL_BACKENDS:
            registry.register(backend_cls)
        return registry


# ---------------------------------------------------------------------------
# Concrete backends
# ---------------------------------------------------------------------------


class PDFBackend(AbstractFormatBackend):
    """PDF extraction: pymupdf (primary) -> pypdf (fallback)."""

    supported_extensions = {".pdf"}

    def extract(self, file_path: str) -> str:
        result = self._try_pymupdf(file_path)
        if result:
            return result
        return self._try_pypdf(file_path)

    def _try_pymupdf(self, file_path: str) -> str:
        try:
            import fitz

            doc = fitz.open(file_path)
            pages = [p.get_text() for p in doc if p.get_text().strip()]
            doc.close()
            return "\n\n".join(pages)
        except ImportError:
            return ""
        except Exception as exc:
            logger.debug("[PDFBackend] pymupdf failed for %s: %s", file_path, exc)
            return ""

    def _try_pypdf(self, file_path: str) -> str:
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            pages = [
                p.extract_text() or ""
                for p in reader.pages
                if (p.extract_text() or "").strip()
            ]
            return "\n\n".join(pages)
        except ImportError:
            self._warn_missing("pymupdf or pypdf", "parsers")
            return ""
        except Exception as exc:
            logger.warning("[PDFBackend] pypdf failed for %s: %s", file_path, exc)
            return ""


class DOCXBackend(AbstractFormatBackend):
    """DOCX extraction: python-docx (paragraphs + tables)."""

    supported_extensions = {".docx", ".dotx", ".docm", ".dotm"}

    def extract(self, file_path: str) -> str:
        try:
            import docx
        except ImportError:
            self._warn_missing("python-docx", "parsers")
            return self._read_plain(file_path)
        try:
            document = docx.Document(file_path)
            blocks = [p.text for p in document.paragraphs if p.text.strip()]
            for table in document.tables:
                blocks.extend(_table_rows(table))
            return "\n".join(blocks)
        except Exception as exc:
            logger.warning("[DOCXBackend] failed for %s: %s", file_path, exc)
            return ""


class DOCBackend(AbstractFormatBackend):
    """Legacy .doc extraction: antiword -> python-docx fallback."""

    supported_extensions = {".doc"}

    def extract(self, file_path: str) -> str:
        result = self._try_antiword(file_path)
        if result:
            return result
        result = self._try_docx(file_path)
        if result:
            return result
        logger.warning(
            "[DOCBackend] .doc extraction failed for %s. Convert to .docx or install antiword.",
            file_path,
        )
        return ""

    def _try_antiword(self, file_path: str) -> str:
        try:
            proc = subprocess.run(
                ["antiword", file_path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if proc.returncode == 0 and proc.stdout.strip():
                return proc.stdout
        except FileNotFoundError:
            logger.debug("[DOCBackend] antiword not found; trying python-docx")
        except subprocess.TimeoutExpired:
            logger.warning("[DOCBackend] antiword timed out for %s", file_path)
        return ""

    def _try_docx(self, file_path: str) -> str:
        try:
            import docx

            document = docx.Document(file_path)
            lines = [p.text for p in document.paragraphs if p.text.strip()]
            return "\n".join(lines)
        except Exception as exc:
            logger.debug(
                "[DOCBackend] python-docx fallback failed for %s: %s", file_path, exc
            )
            return ""


class XLSXBackend(AbstractFormatBackend):
    """Excel extraction: openpyxl (multi-sheet)."""

    supported_extensions = {".xlsx", ".xls", ".xlsm"}

    def extract(self, file_path: str) -> str:
        try:
            import openpyxl
        except ImportError:
            self._warn_missing("openpyxl", "parsers")
            return self._read_plain(file_path)
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        try:
            lines: List[str] = []
            for sheet in wb.worksheets:
                lines.append(f"[Sheet: {sheet.title}]")
                lines.extend(_xlsx_sheet_rows(sheet))
            return "\n".join(lines)
        finally:
            wb.close()


class PPTXBackend(AbstractFormatBackend):
    """PowerPoint extraction: python-pptx (text + tables)."""

    supported_extensions = {
        ".pptx",
        ".potx",
        ".ppsx",
        ".pptm",
        ".potm",
        ".ppsm",
        ".ppt",
    }

    def extract(self, file_path: str) -> str:
        try:
            from pptx import Presentation
            from pptx.shapes.graphfrm import GraphicFrame
        except ImportError:
            self._warn_missing("python-pptx", "parsers")
            return ""
        prs = Presentation(file_path)
        lines: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"[Slide {i}]")
            lines.extend(_slide_text(slide, GraphicFrame))
        return "\n".join(lines)


class RTFBackend(AbstractFormatBackend):
    """RTF extraction: striprtf -> plain-text fallback."""

    supported_extensions = {".rtf"}

    def extract(self, file_path: str) -> str:
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            self._warn_missing("striprtf", "parsers")
            return self._read_plain(file_path)
        with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
            return rtf_to_text(fh.read())


class HTMLBackend(AbstractFormatBackend):
    """HTML/XHTML extraction: stdlib html.parser (zero deps)."""

    supported_extensions = {".html", ".htm", ".xhtml"}

    def extract(self, file_path: str) -> str:
        return extract_html(file_path)


class VTTBackend(AbstractFormatBackend):
    """WebVTT subtitle extraction: internal regex parser."""

    supported_extensions = {".vtt"}

    def extract(self, file_path: str) -> str:
        return extract_vtt(file_path)


class LaTeXBackend(AbstractFormatBackend):
    """LaTeX extraction: internal regex stripper."""

    supported_extensions = {".tex", ".latex"}

    def extract(self, file_path: str) -> str:
        return extract_latex(file_path)


class CSVBackend(AbstractFormatBackend):
    """CSV extraction: stdlib csv.Sniffer auto-delimiter detection."""

    supported_extensions = {".csv"}

    def extract(self, file_path: str) -> str:
        return extract_csv(file_path)


class ImageBackend(AbstractFormatBackend):
    """Image OCR extraction: pytesseract ([vision] extra)."""

    supported_extensions = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".webp"}

    def extract(self, file_path: str) -> str:
        return extract_image(file_path)


class PlainTextBackend(AbstractFormatBackend):
    """Plain UTF-8 read for .txt, .md, and unknown text formats."""

    supported_extensions = {
        ".md",
        ".txt",
        ".text",
        ".qmd",
        ".rmd",
    }

    def extract(self, file_path: str) -> str:
        return self._read_plain(file_path)


# ---------------------------------------------------------------------------
# Ordered registration list (determines dispatch priority)
# ---------------------------------------------------------------------------

_ALL_BACKENDS: List[Type[AbstractFormatBackend]] = [
    PDFBackend,
    DOCXBackend,
    DOCBackend,
    XLSXBackend,
    PPTXBackend,
    RTFBackend,
    HTMLBackend,
    VTTBackend,
    LaTeXBackend,
    CSVBackend,
    ImageBackend,
    PlainTextBackend,
]


# ---------------------------------------------------------------------------
# Private helpers shared by multiple backends
# ---------------------------------------------------------------------------


def _table_rows(table: object) -> List[str]:
    """Extract tab-separated rows from a docx/pptx table object."""
    rows: List[str] = []
    for row in table.rows:  # type: ignore[attr-defined]
        row_text = "\t".join(c.text.strip() for c in row.cells if c.text.strip())
        if row_text:
            rows.append(row_text)
    return rows


def _xlsx_sheet_rows(sheet: object) -> List[str]:
    """Extract non-empty rows from an openpyxl worksheet."""
    rows: List[str] = []
    for row in sheet.iter_rows(values_only=True):  # type: ignore[attr-defined]
        row_text = "\t".join("" if c is None else str(c) for c in row)
        if row_text.strip():
            rows.append(row_text)
    return rows


def _slide_text(slide: object, graphic_frame_cls: type) -> List[str]:
    """Extract text from all shapes on a single PPTX slide."""
    parts: List[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if isinstance(shape, graphic_frame_cls) and shape.has_table:
            parts.extend(_table_rows(shape.table))
        elif hasattr(shape, "text") and shape.text.strip():
            parts.append(shape.text.strip())
    return parts
