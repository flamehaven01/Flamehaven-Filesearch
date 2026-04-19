"""
File format text extractors for FLAMEHAVEN FileSearch.

All parsers are implemented internally — no external document-AI framework
dependency. Optional heavy libs (pymupdf, python-docx, etc.) are guarded
with ImportError fallbacks and installed via pyproject extras.

Extraction stack (tried in order until non-empty result):
  PDF   : pymupdf  ->  pypdf  ->  plain-text fallback
  DOCX  : python-docx (tables included)
  DOC   : antiword  ->  python-docx  ->  empty
  XLSX  : openpyxl
  PPTX  : python-pptx (tables included)
  RTF   : striprtf
  HTML  : internal html.parser extractor (stdlib, no deps)
  VTT   : internal WebVTT regex parser   (stdlib, no deps)
  LaTeX : internal regex stripper        (stdlib, no deps)
  CSV   : internal csv.Sniffer parser    (stdlib, no deps)
  Image : pytesseract OCR               (optional [vision] extra)
  TXT/MD: plain UTF-8 read

All functions return plain UTF-8 text suitable for embedding and RAG.
For RAG chunking use engine.text_chunker.chunk_text().
"""

import logging
from pathlib import Path
from typing import List

from .format_parsers import (
    extract_csv,
    extract_html,
    extract_image,
    extract_latex,
    extract_vtt,
)

logger = logging.getLogger(__name__)

# Image extensions supported via OCR ([vision] extra)
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".webp"}

SUPPORTED_EXTENSIONS = {
    # Documents (require [parsers] extra)
    ".pdf", ".docx", ".dotx", ".docm", ".dotm",
    ".doc",
    ".xlsx", ".xlsm",
    ".pptx", ".potx", ".ppsx", ".pptm", ".potm", ".ppsm",
    ".rtf",
    # Markup / text (stdlib, zero deps)
    ".html", ".htm", ".xhtml",
    ".md", ".txt", ".text", ".qmd", ".rmd",
    ".csv",
    ".vtt",
    ".tex", ".latex",
    # Images (require [vision] extra)
    ".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".webp",
}


def extract_text(file_path: str, use_cache: bool = False) -> str:
    """Extract plain UTF-8 text from a file based on its extension.

    Falls back to plain-text read for unrecognised formats.
    Returns empty string on unrecoverable errors.

    Args:
        file_path: Path to the source file.
        use_cache: When True, check parse_cache before extracting and store
                   the result after a successful parse. Cache is invalidated
                   automatically via file mtime — no manual flush needed for
                   normal file updates.
    """
    if use_cache:
        from .parse_cache import get as _cache_get, put as _cache_put
        cached = _cache_get(file_path)
        if cached is not None:
            return cached

    result = _dispatch_extract(file_path)

    if use_cache:
        from .parse_cache import put as _cache_put
        _cache_put(file_path, result)

    return result


def _dispatch_extract(file_path: str) -> str:
    """Route extraction to the appropriate parser by file extension."""
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        if ext in {".docx", ".dotx", ".docm", ".dotm"}:
            return _extract_docx(file_path)
        if ext == ".doc":
            return _extract_doc(file_path)
        if ext in {".xlsx", ".xls", ".xlsm"}:
            return _extract_xlsx(file_path)
        if ext in {".pptx", ".potx", ".ppsx", ".pptm", ".potm", ".ppsm", ".ppt"}:
            return _extract_pptx(file_path)
        if ext == ".rtf":
            return _extract_rtf(file_path)
        if ext in {".html", ".htm", ".xhtml"}:
            return extract_html(file_path)
        if ext == ".vtt":
            return extract_vtt(file_path)
        if ext in {".tex", ".latex"}:
            return extract_latex(file_path)
        if ext == ".csv":
            return extract_csv(file_path)
        if ext in _IMAGE_EXTS:
            return extract_image(file_path)
        return _read_plain(file_path)
    except Exception as exc:
        logger.warning("[FileParser] Failed to extract %s: %s", file_path, exc)
        return ""


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF via pymupdf (primary) or pypdf (fallback)."""
    try:
        import fitz  # pymupdf

        doc = fitz.open(file_path)
        pages = [page.get_text() for page in doc if page.get_text().strip()]
        doc.close()
        if pages:
            return "\n\n".join(pages)
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("[FileParser] pymupdf failed for %s: %s", file_path, exc)

    try:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        pages = [p.extract_text() or "" for p in reader.pages if (p.extract_text() or "").strip()]
        return "\n\n".join(pages)
    except ImportError:
        logger.warning(
            "[FileParser] PDF extraction requires pymupdf or pypdf. "
            "Install: pip install flamehaven-filesearch[parsers]"
        )
        return ""
    except Exception as exc:
        logger.warning("[FileParser] pypdf failed for %s: %s", file_path, exc)
        return ""


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def _extract_docx(file_path: str) -> str:
    """Extract text from DOCX via python-docx, including tables."""
    try:
        import docx
    except ImportError:
        logger.warning(
            "[FileParser] python-docx not installed. "
            "Install: pip install flamehaven-filesearch[parsers]"
        )
        return _read_plain(file_path)

    try:
        document = docx.Document(file_path)
        blocks = [para.text for para in document.paragraphs if para.text.strip()]
        for table in document.tables:
            blocks.extend(_extract_table_rows(table))
        return "\n".join(blocks)
    except Exception as exc:
        logger.warning("[FileParser] DOCX extraction failed for %s: %s", file_path, exc)
        return ""


# ---------------------------------------------------------------------------
# DOC (legacy binary Word)
# ---------------------------------------------------------------------------


def _extract_doc(file_path: str) -> str:
    """Extract text from legacy .doc via antiword or python-docx fallback."""
    import subprocess

    try:
        result = subprocess.run(
            ["antiword", file_path],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        logger.debug("[FileParser] antiword not found; trying python-docx fallback")
    except subprocess.TimeoutExpired:
        logger.warning("[FileParser] antiword timed out for %s", file_path)

    try:
        import docx

        document = docx.Document(file_path)
        lines = [p.text for p in document.paragraphs if p.text.strip()]
        if lines:
            return "\n".join(lines)
    except Exception as exc:
        logger.debug("[FileParser] python-docx fallback failed for %s: %s", file_path, exc)

    logger.warning(
        "[FileParser] .doc extraction failed for %s. "
        "Convert to .docx or install antiword.",
        file_path,
    )
    return ""


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def _extract_xlsx(file_path: str) -> str:
    """Extract text from Excel workbooks via openpyxl."""
    try:
        import openpyxl
    except ImportError:
        logger.warning(
            "[FileParser] openpyxl not installed; "
            "run: pip install flamehaven-filesearch[parsers]"
        )
        return _read_plain(file_path)

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    try:
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"[Sheet: {sheet.title}]")
            lines.extend(_extract_xlsx_sheet(sheet))
        return "\n".join(lines)
    finally:
        wb.close()


def _extract_xlsx_sheet(sheet: object) -> List[str]:
    """Extract non-empty tab-separated rows from an openpyxl worksheet."""
    rows = []
    for row in sheet.iter_rows(values_only=True):  # type: ignore[attr-defined]
        row_text = "\t".join("" if c is None else str(c) for c in row)
        if row_text.strip():
            rows.append(row_text)
    return rows


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def _extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint presentations via python-pptx."""
    try:
        from pptx import Presentation
        from pptx.shapes.graphfrm import GraphicFrame
    except ImportError:
        logger.warning(
            "[FileParser] python-pptx not installed; "
            "run: pip install flamehaven-filesearch[parsers]"
        )
        return ""

    prs = Presentation(file_path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {i}]")
        lines.extend(_extract_slide_text(slide, GraphicFrame))
    return "\n".join(lines)


def _extract_slide_text(slide: object, graphic_frame_cls: type) -> List[str]:
    """Extract text from all shapes on a single PPTX slide."""
    parts: List[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if isinstance(shape, graphic_frame_cls) and shape.has_table:
            parts.extend(_extract_table_rows(shape.table))
        elif hasattr(shape, "text") and shape.text.strip():
            parts.append(shape.text.strip())
    return parts


def _extract_table_rows(table: object) -> List[str]:
    """Extract tab-separated rows from a pptx/docx table object."""
    rows = []
    for row in table.rows:  # type: ignore[attr-defined]
        row_text = "\t".join(
            cell.text.strip() for cell in row.cells if cell.text.strip()
        )
        if row_text:
            rows.append(row_text)
    return rows


# ---------------------------------------------------------------------------
# RTF
# ---------------------------------------------------------------------------


def _extract_rtf(file_path: str) -> str:
    """Extract plain text from RTF documents via striprtf."""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        logger.warning(
            "[FileParser] striprtf not installed; "
            "run: pip install flamehaven-filesearch[parsers]"
        )
        return _read_plain(file_path)

    with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
        return rtf_to_text(fh.read())


# ---------------------------------------------------------------------------
# Plain text fallback
# ---------------------------------------------------------------------------


def _read_plain(file_path: str) -> str:
    """Read file as plain UTF-8 text (last-resort fallback)."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
            return fh.read()
    except OSError:
        return ""
